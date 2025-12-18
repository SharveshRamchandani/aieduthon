from datetime import datetime
from io import BytesIO
from pathlib import Path
from typing import Dict, List, Optional, Tuple
import logging
import base64
import io
import os

import requests
from bson.objectid import ObjectId
from gridfs import GridFS
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image

from ai_db import get_ai_db

logger = logging.getLogger(__name__)


class PPTExporter:
	def __init__(self):
		self.db = get_ai_db()
		self.slides_collection = self.db["slides"]

	def _capture_paragraph_style(self, paragraph) -> Dict[str, Optional[object]]:
		"""
		Snapshot the key visual properties from a template paragraph so we can
		reuse them when injecting our own text. This keeps font, size, color,
		alignment and spacing consistent with the original slide design.
		"""
		if paragraph is None:
			return {}

		style: Dict[str, Optional[object]] = {}
		try:
			font = paragraph.font
			style["font_name"] = font.name
			style["font_size"] = font.size
			style["bold"] = font.bold
			style["italic"] = font.italic
			# color can be None or theme-based; guard for rgb
			try:
				if font.color is not None and getattr(font.color, "rgb", None):
					style["font_color"] = font.color.rgb
				else:
					style["font_color"] = None
			except Exception:
				style["font_color"] = None

			style["level"] = paragraph.level
			style["alignment"] = paragraph.alignment
			style["space_after"] = paragraph.space_after
			style["space_before"] = paragraph.space_before
		except Exception:
			# Best-effort only; if anything fails we just return what we captured
			pass

		return style

	def _apply_paragraph_style(self, paragraph, style: Dict[str, Optional[object]]):
		"""Apply a previously captured style dictionary to a paragraph."""
		if not style:
			return
		try:
			font = paragraph.font
			if "font_name" in style and style["font_name"] is not None:
				font.name = style["font_name"]
			if "font_size" in style and style["font_size"] is not None:
				font.size = style["font_size"]
			if "bold" in style and style["bold"] is not None:
				font.bold = style["bold"]
			if "italic" in style and style["italic"] is not None:
				font.italic = style["italic"]
			if "font_color" in style and style["font_color"] is not None:
				try:
					font.color.rgb = style["font_color"]
				except Exception:
					pass

			if "level" in style and style["level"] is not None:
				paragraph.level = style["level"]
			if "alignment" in style and style["alignment"] is not None:
				paragraph.alignment = style["alignment"]
			if "space_after" in style and style["space_after"] is not None:
				paragraph.space_after = style["space_after"]
			if "space_before" in style and style["space_before"] is not None:
				paragraph.space_before = style["space_before"]
		except Exception:
			# Styling issues should never break export
			pass

	def _calculate_available_space_for_image(self, slide, text_frame) -> Dict[str, float]:
		"""
		Calculate available space on slide for image placement based on text content.
		
		Args:
			slide: PowerPoint slide object
			text_frame: Text frame object containing slide content
			
		Returns:
			Dict with available space dimensions
		"""
		# Standard slide dimensions (10x7.5 inches for 16:9)
		SLIDE_WIDTH = 10.0
		SLIDE_HEIGHT = 7.5
		MARGIN = 0.5
		
		# Estimate content area from text frame position
		try:
			# Get parent shape to access position
			parent_shape = text_frame._parent  # Access parent shape
			text_left = parent_shape.left / 914400.0  # Convert EMU to inches
			text_top = parent_shape.top / 914400.0
			text_width = parent_shape.width / 914400.0
			text_height = parent_shape.height / 914400.0
		except Exception:
			# Fallback: estimate content area
			text_left = 1.0
			text_top = 1.5
			text_width = 5.5
			text_height = 5.5
		
		# Calculate available space for image (right side)
		available_left = text_left + text_width + 0.2
		available_top = text_top
		available_width = SLIDE_WIDTH - available_left - MARGIN
		available_height = SLIDE_HEIGHT - available_top - MARGIN
		
		return {
			"left": max(MARGIN, available_left),
			"top": max(MARGIN, available_top),
			"width": max(2.0, available_width),  # Minimum 2 inches
			"height": max(2.0, available_height),  # Minimum 2 inches
		}
	
	def _calculate_optimal_image_size(
		self,
		image_path: str,
		available_space: Dict[str, float],
		max_size_ratio: float = 0.85,
		min_size_ratio: float = 0.4
	) -> Dict[str, float]:
		"""
		Calculate optimal image size based on available space and image dimensions.
		
		Args:
			image_path: Path to image file
			available_space: Dict with available space dimensions
			max_size_ratio: Maximum ratio of available space to use (0.0-1.0)
			min_size_ratio: Minimum ratio of available space to use (0.0-1.0)
			
		Returns:
			Dict with optimal width, height, left, top positions
		"""
		try:
			with Image.open(image_path) as image:
				img_width_px, img_height_px = image.size
				img_ratio = img_width_px / img_height_px if img_height_px > 0 else 1.0
		except Exception:
			# Fallback if image can't be opened
			img_ratio = 1.0
		
		avail_w = available_space["width"]
		avail_h = available_space["height"]
		avail_ratio = avail_w / avail_h if avail_h > 0 else 1.0
		
		# Calculate optimal size maintaining aspect ratio
		if img_ratio > avail_ratio:
			# Image is wider than available space
			optimal_width = avail_w * max_size_ratio
			optimal_height = optimal_width / img_ratio
			# Ensure minimum size
			if optimal_height < avail_h * min_size_ratio:
				optimal_height = avail_h * min_size_ratio
				optimal_width = optimal_height * img_ratio
		else:
			# Image is taller than available space
			optimal_height = avail_h * max_size_ratio
			optimal_width = optimal_height * img_ratio
			# Ensure minimum size
			if optimal_width < avail_w * min_size_ratio:
				optimal_width = avail_w * min_size_ratio
				optimal_height = optimal_width / img_ratio
		
		# Center the image in available space
		left = available_space["left"] + (avail_w - optimal_width) / 2
		top = available_space["top"] + (avail_h - optimal_height) / 2
		
		return {
			"left": max(available_space["left"], left),
			"top": max(available_space["top"], top),
			"width": min(optimal_width, avail_w),
			"height": min(optimal_height, avail_h)
		}

	def export_deck(self, deck_id: str, output_dir: str = "..\\..\\out", user_name: str = "user", export_style: str = "template") -> str:
		"""Export slide deck to disk (legacy behavior).
		
		export_style:
			- "template": use the configured PPT template (default)
			- "preview": ignore template and render slides in the dark preview style
		"""
		prs, fname = self._build_presentation(deck_id, user_name, export_style=export_style)
		out_dir = Path(output_dir).resolve()
		out_dir.mkdir(parents=True, exist_ok=True)
		output_path = out_dir / fname
		prs.save(str(output_path))
		return str(output_path)

	def export_deck_to_bytes(
		self,
		deck_id: str,
		save_to_db: bool = True,
		user_name: str = "user",
		format_type: str = "pptx",
		export_style: str = "template",
	) -> Tuple[bytes, str]:
		"""Return PPTX or PDF content as bytes plus suggested filename. Optionally save to database.
		
		Args:
			deck_id: Deck ID
			save_to_db: Whether to save to database
			user_name: User name for filename
			format_type: "pptx" or "pdf"
			export_style: "template" or "preview"
		"""
		if format_type == "pdf":
			# For now, PDF export always uses the template-based presentation
			# to preserve structure; we can later add a dedicated preview PDF
			# layout if needed.
			return self._export_to_pdf(deck_id, user_name, export_style=export_style)
		
		prs, fname = self._build_presentation(deck_id, user_name, export_style=export_style)
		buffer = BytesIO()
		prs.save(buffer)
		ppt_bytes = buffer.getvalue()
		
		# Save PPT to database if requested
		if save_to_db:
			self._save_ppt_to_db(deck_id, ppt_bytes, fname)
		
		return ppt_bytes, fname
	
	def _export_to_pdf(self, deck_id: str, user_name: str = "user", export_style: str = "template") -> Tuple[bytes, str]:
		"""Convert PPTX to PDF preserving template and formatting.
		
		Order of converters:
		1) Microsoft PowerPoint COM automation (Windows only, if PowerPoint installed)
		2) LibreOffice headless (if available locally or via SOFFICE_PATH)
		3) ReportLab fallback (loses template/layout fidelity)
		"""
		import re
		import tempfile
		import subprocess
		import shutil
		import platform
		
		try:
			object_id = ObjectId(deck_id)
		except Exception:
			raise ValueError("Invalid deck_id. Must be a Mongo ObjectId hex string.")

		deck = self.slides_collection.find_one({"_id": object_id})
		if not deck:
			raise FileNotFoundError("Slide deck not found")
		
		title = deck.get("title", "AI Presentation")
		
		# Generate PPTX first (in memory)
		prs, pptx_filename = self._build_presentation(deck_id, user_name, export_style=export_style)
		
		# Generate filename: Topic_user_date.pdf
		safe_title = re.sub(r'[^\w\s-]', '', title)[:50].strip().replace(' ', '_')
		if not safe_title:
			safe_title = "Presentation"
		
		safe_user = re.sub(r'[^\w\s-]', '', user_name)[:30].strip().replace(' ', '_')
		if not safe_user:
			safe_user = "user"
		
		date_str = datetime.utcnow().strftime('%Y-%m-%d')
		filename = f"{safe_title}_{safe_user}_{date_str}.pdf"
		
		# 1) Try Microsoft PowerPoint COM automation on Windows
		if platform.system() == "Windows":
			try:
				return self._convert_pptx_to_pdf_powerpoint(prs, filename)
			except Exception as ppt_com_error:
				logger.warning(f"PowerPoint COM conversion failed: {ppt_com_error}. Trying LibreOffice next.")

		# 2) Try LibreOffice conversion (preserves template and formatting)
		try:
			return self._convert_pptx_to_pdf_libreoffice(prs, filename)
		except Exception as libreoffice_error:
			logger.warning(f"LibreOffice conversion failed: {libreoffice_error}. Falling back to reportlab (template/formatting may be lost).")

		# 3) Fallback to reportlab (loses template but at least produces PDF)
		return self._convert_pptx_to_pdf_reportlab(prs, title, user_name, filename)

	def _convert_pptx_to_pdf_powerpoint(self, prs: Presentation, filename: str) -> Tuple[bytes, str]:
		"""Use installed Microsoft PowerPoint via COM to export PDF (Windows only)."""
		import tempfile
		import platform

		if platform.system() != "Windows":
			raise RuntimeError("PowerPoint COM export is only available on Windows.")

		try:
			import win32com.client  # type: ignore
		except Exception as e:
			raise RuntimeError(f"pywin32 not available for PowerPoint COM export: {e}")

		with tempfile.TemporaryDirectory() as temp_dir:
			pptx_path = Path(temp_dir) / "presentation.pptx"
			pdf_path = Path(temp_dir) / "presentation.pdf"

			# Save PPTX to temp
			prs.save(str(pptx_path))

			powerpoint = None
			try:
				powerpoint = win32com.client.Dispatch("PowerPoint.Application")
				powerpoint.Visible = 0
				presentation = powerpoint.Presentations.Open(str(pptx_path), WithWindow=False)
				# 32 is the PDF format in PowerPoint SaveAs
				presentation.SaveAs(str(pdf_path), 32)
				presentation.Close()
			finally:
				if powerpoint:
					powerpoint.Quit()

			if not pdf_path.exists():
				raise RuntimeError("PowerPoint did not produce PDF file")

			with open(pdf_path, "rb") as f:
				pdf_bytes = f.read()

		return pdf_bytes, filename
	
	def _convert_pptx_to_pdf_libreoffice(self, prs: Presentation, filename: str) -> Tuple[bytes, str]:
		"""Convert PPTX to PDF using LibreOffice headless mode - preserves all formatting."""
		import tempfile
		import subprocess
		import shutil
		import platform
		
		# Check if LibreOffice is available. Allow overriding via SOFFICE_PATH.
		env_soffice = os.getenv("SOFFICE_PATH")
		if env_soffice and Path(env_soffice).exists():
			soffice = env_soffice
		elif platform.system() == "Windows":
			# Common LibreOffice paths on Windows
			libreoffice_paths = [
				r"C:\Program Files\LibreOffice\program\soffice.exe",
				r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
			]
			soffice = None
			for path in libreoffice_paths:
				if Path(path).exists():
					soffice = path
					break
			if not soffice:
				# Try to find in PATH
				soffice = shutil.which("soffice")
		else:
			# Linux/Mac - try to find in PATH
			soffice = shutil.which("soffice") or shutil.which("libreoffice")
		
		if not soffice:
			raise RuntimeError("LibreOffice not found. Install LibreOffice for PDF export with template preservation.")
		
		# Save PPTX to temporary file
		with tempfile.TemporaryDirectory() as temp_dir:
			pptx_path = Path(temp_dir) / "presentation.pptx"
			pdf_path = Path(temp_dir) / "presentation.pdf"
			
			# Save PPTX
			prs.save(str(pptx_path))
			
			# Convert using LibreOffice headless mode
			# --headless: Run without GUI
			# --convert-to pdf: Convert to PDF
			# --outdir: Output directory
			cmd = [
				soffice,
				"--headless",
				"--convert-to", "pdf",
				"--outdir", str(temp_dir),
				str(pptx_path)
			]
			
			try:
				result = subprocess.run(
					cmd,
					capture_output=True,
					text=True,
					timeout=60,  # 60 second timeout
					check=True
				)
			except subprocess.TimeoutExpired:
				raise RuntimeError("LibreOffice conversion timed out")
			except subprocess.CalledProcessError as e:
				raise RuntimeError(f"LibreOffice conversion failed: {e.stderr}")
			
			# Check if PDF was created
			if not pdf_path.exists():
				raise RuntimeError("LibreOffice did not produce PDF file")
			
			# Read PDF bytes
			with open(pdf_path, 'rb') as f:
				pdf_bytes = f.read()
			
			return pdf_bytes, filename
	
	def _convert_pptx_to_pdf_reportlab(self, prs: Presentation, title: str, user_name: str, filename: str) -> Tuple[bytes, str]:
		"""Fallback PDF conversion using reportlab - loses template but produces PDF."""
		try:
			from reportlab.lib.pagesizes import letter, A4
			from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
			from reportlab.lib.units import inch
			from reportlab.lib import colors
			from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak, Image as RLImage
			from reportlab.lib.enums import TA_CENTER, TA_LEFT
		except ImportError:
			raise ImportError("reportlab is required for PDF export fallback. Install with: pip install reportlab")
		
		import tempfile
		
		# Create PDF buffer
		pdf_buffer = BytesIO()
		
		# Use letter size (8.5 x 11 inches) for slides
		doc = SimpleDocTemplate(
			pdf_buffer,
			pagesize=letter,
			rightMargin=0.5*inch,
			leftMargin=0.5*inch,
			topMargin=0.5*inch,
			bottomMargin=0.5*inch
		)
		
		styles = getSampleStyleSheet()
		
		# Custom styles for slides
		slide_title_style = ParagraphStyle(
			'SlideHeading',
			parent=styles['Heading2'],
			fontSize=24,
			textColor=colors.HexColor('#2c3e50'),
			spaceAfter=12,
			spaceBefore=10,
			fontName='Helvetica-Bold'
		)
		
		body_style = ParagraphStyle(
			'SlideBody',
			parent=styles['Normal'],
			fontSize=14,
			leading=18,
			spaceAfter=10,
			textColor=colors.HexColor('#333333'),
			fontName='Helvetica'
		)
		
		story = []
		temp_image_files = []  # Keep track of temp files to clean up after PDF is built
		
		# Process each slide
		for slide_idx, slide in enumerate(prs.slides):
			# Extract slide title
			slide_title = ""
			try:
				if slide.shapes.title:
					slide_title = slide.shapes.title.text.strip()
			except:
				pass
			
			# Add slide title if present
			if slide_title:
				story.append(Paragraph(slide_title, slide_title_style))
				story.append(Spacer(1, 0.2*inch))
			
			# Extract text content from all shapes
			text_content = []
			images = []
			
			for shape in slide.shapes:
				# Extract text from text boxes and placeholders
				if hasattr(shape, "text") and shape.text:
					text = shape.text.strip()
					if text and text != slide_title:  # Don't duplicate title
						text_content.append(text)
				
				# Extract images
				if hasattr(shape, "image"):
					try:
						image = shape.image
						image_bytes = image.blob
						# Save to temp file for reportlab
						with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as tmp_img:
							tmp_img.write(image_bytes)
							img_path = tmp_img.name
							temp_image_files.append(img_path)  # Track for cleanup
							images.append((img_path, shape))
					except Exception as e:
						logger.debug(f"Failed to extract image from slide {slide_idx}: {e}")
			
			# Add text content
			for text in text_content:
				# Clean HTML-like tags and format
				clean_text = text.replace('\n', '<br/>')
				story.append(Paragraph(clean_text, body_style))
				story.append(Spacer(1, 0.1*inch))
			
			# Add images (if any)
			for img_path, shape in images:
				try:
					# Calculate image size (maintain aspect ratio, max width 7 inches)
					from PIL import Image as PILImage
					with PILImage.open(img_path) as pil_img:
						img_width_px, img_height_px = pil_img.size
						aspect = img_width_px / img_height_px
					
					# Get shape dimensions if available, otherwise use defaults
					try:
						shape_width = shape.width
						shape_height = shape.height
						# Convert EMU to inches (1 inch = 914400 EMU)
						shape_width_inch = shape_width / 914400
						shape_height_inch = shape_height / 914400
					except:
						shape_width_inch = 5.0
						shape_height_inch = 3.75
					
					# Use shape dimensions or calculate from image pixels
					max_width = 7 * inch
					max_height = 5 * inch
					
					# Use shape dimensions if reasonable, otherwise calculate from image
					if shape_width_inch > 0 and shape_height_inch > 0:
						width = min(shape_width_inch * inch, max_width)
						height = min(shape_height_inch * inch, max_height)
					else:
						# Calculate from image pixels (assuming 96 DPI)
						dpi = 96
						width_inch = img_width_px / dpi
						height_inch = img_height_px / dpi
						
						if width_inch > max_width / inch:
							width = max_width
							height = width / aspect
						elif height_inch > max_height / inch:
							height = max_height
							width = height * aspect
						else:
							width = width_inch * inch
							height = height_inch * inch
					
					story.append(RLImage(img_path, width=width, height=height))
					story.append(Spacer(1, 0.2*inch))
				except Exception as e:
					logger.debug(f"Failed to add image to PDF: {e}")
			
			# Add page break after each slide (except last)
			if slide_idx < len(prs.slides) - 1:
				story.append(PageBreak())
		
		# Build PDF (this is when ReportLab reads the image files)
		try:
			doc.build(story)
			pdf_bytes = pdf_buffer.getvalue()
		finally:
			# Clean up temp image files AFTER PDF is built
			for img_path in temp_image_files:
				try:
					if Path(img_path).exists():
						Path(img_path).unlink()
				except Exception as e:
					logger.debug(f"Failed to delete temp image file {img_path}: {e}")
		
		return pdf_bytes, filename
	
	def _save_ppt_to_db(self, deck_id: str, ppt_bytes: bytes, filename: str):
		"""Save generated PPT file to database"""
		try:
			object_id = ObjectId(deck_id)
		except Exception:
			logger.warning(f"Invalid deck_id for PPT storage: {deck_id}")
			return
		
		# Update the deck document with PPT file info
		import base64
		ppt_base64 = base64.b64encode(ppt_bytes).decode("utf-8")
		
		update_data = {
			"ppt_file": ppt_base64,  # Store as base64 string
			"ppt_filename": filename,
			"ppt_generated_at": datetime.utcnow(),
			"ppt_size_bytes": len(ppt_bytes)
		}
		
		self.slides_collection.update_one(
			{"_id": object_id},
			{"$set": update_data}
		)
		logger.info(f"Saved PPT file to database for deck {deck_id}")

	def _build_preview_presentation(self, deck_id: str, user_name: str = "user") -> Tuple[Presentation, str]:
		"""Build a PPT that mimics the web preview: dark background, white text, full-bleed image."""
		try:
			object_id = ObjectId(deck_id)
		except Exception:
			raise ValueError("Invalid deck_id. Must be a Mongo ObjectId hex string.")

		deck = self.slides_collection.find_one({"_id": object_id})
		if not deck:
			raise FileNotFoundError("Slide deck not found")

		title = deck.get("title", "Presentation")
		sections: List[str] = deck.get("sections", [])
		bullets: List[List[str]] = deck.get("bullets", [])
		expanded_content: List[str] = deck.get("expanded_content", [])
		generated_notes = deck.get("generated_notes", [])
		media_refs = deck.get("media_refs", [])

		prs = Presentation()
		# Try to use a blank layout for maximum control
		blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]

		from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE

		for idx, section in enumerate(sections):
			slide = prs.slides.add_slide(blank_layout)

			# Dark background similar to preview (slate-900)
			bg_fill = slide.background.fill
			bg_fill.solid()
			bg_fill.fore_color.rgb = RGBColor(15, 23, 42)

			# Full-bleed background image from first media ref if available
			text_rgb = RGBColor(255, 255, 255)  # default to light text
			slide_media = media_refs[idx] if idx < len(media_refs) else []
			url = slide_media[0] if slide_media and isinstance(slide_media[0], str) else None
			if url:
				try:
					resp = requests.get(url, timeout=15)
					resp.raise_for_status()
					tmp_path = Path(f"_preview_bg_{idx}.png")
					tmp_path.write_bytes(resp.content)
					try:
						# Add image as full-bleed background
						slide.shapes.add_picture(str(tmp_path), 0, 0, width=prs.slide_width, height=prs.slide_height)

						# Sample brightness to choose light/dark text automatically
						try:
							with Image.open(str(tmp_path)) as im:
								# Downscale for speed, convert to grayscale, compute mean brightness
								im_small = im.resize((64, 64)).convert("L")
								avg = sum(im_small.getdata()) / float(len(im_small.getdata()))
								# If background is bright, use darker text; otherwise keep light text
								if avg > 150:
									# darker slate text
									text_rgb = RGBColor(31, 41, 55)  # slate-800
								else:
									text_rgb = RGBColor(248, 250, 252)  # slate-50
						except Exception as sample_err:
							logger.debug(f"Failed to sample background brightness for slide {idx}: {sample_err}")
					finally:
						try:
							tmp_path.unlink()
						except OSError:
							pass
				except Exception as e:
					logger.debug(f"Failed to add preview background image for slide {idx}: {e}")

			# Derive main body text from generated notes / expanded content / bullets
			expanded_text = generated_notes[idx] if idx < len(generated_notes) else ""
			if not expanded_text and idx < len(expanded_content):
				expanded_text = expanded_content[idx]
			if not expanded_text and idx < len(bullets):
				expanded_text = "\n".join(bullets[idx])

			# Title textbox (top center, slightly inset from sides)
			title_box = slide.shapes.add_textbox(
				Inches(0.8),
				Inches(0.8),
				prs.slide_width - Inches(1.6),
				Inches(2.0),
			)
			title_tf = title_box.text_frame
			title_tf.clear()
			try:
				title_tf.word_wrap = True
				title_tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
			except Exception:
				pass
			p_title = title_tf.paragraphs[0]
			p_title.text = section or title
			p_title.alignment = PP_ALIGN.CENTER
			title_font = p_title.font
			title_font.name = "Arial"
			title_font.size = Pt(40)
			title_font.bold = True
			title_font.color.rgb = text_rgb

			# Body textbox (centered, text color adapts to background)
			# Body box starts a bit lower to avoid overlapping the title and has
			# generous side margins so long lines wrap instead of running edge to edge.
			body_box = slide.shapes.add_textbox(
				Inches(0.9),
				Inches(3.0),
				prs.slide_width - Inches(1.8),
				prs.slide_height - Inches(3.6),
			)
			body_tf = body_box.text_frame
			body_tf.clear()
			# Let PowerPoint try to shrink text to fit the shape and wrap lines
			try:
				body_tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
				body_tf.word_wrap = True
				body_tf.margin_left = Inches(0.2)
				body_tf.margin_right = Inches(0.2)
			except Exception:
				pass
			p_body = body_tf.paragraphs[0]
			p_body.alignment = PP_ALIGN.CENTER
			body_font = p_body.font
			body_font.name = "Arial"

			# Heuristic font size based on total text length so long paragraphs
			# shrink automatically and stay inside the frame.
			total_len = len(str(expanded_text))
			if total_len <= 500:
				base_size = 24
			elif total_len <= 900:
				base_size = 18
			elif total_len <= 1300:
				base_size = 16
			else:
				base_size = 14
			body_font.size = Pt(base_size)
			body_font.color.rgb = text_rgb

			if expanded_text:
				for i, line in enumerate(str(expanded_text).split("\n")):
					text = line.strip()
					if not text:
						continue
					if i == 0:
						p = p_body
					else:
						p = body_tf.add_paragraph()
						p.alignment = PP_ALIGN.CENTER
						p.font.name = body_font.name
						p.font.size = body_font.size
						p.font.color.rgb = body_font.color.rgb
					p.text = text

		# Add a minimalistic "Thank you" slide at the end in preview mode
		if sections:
			last_media = media_refs[len(sections) - 1] if len(media_refs) >= len(sections) else []
			last_url = last_media[0] if last_media and isinstance(last_media[0], str) else None

			slide = prs.slides.add_slide(blank_layout)

			# Dark background
			bg_fill = slide.background.fill
			bg_fill.solid()
			bg_fill.fore_color.rgb = RGBColor(15, 23, 42)

			text_rgb = RGBColor(248, 250, 252)

			# Use last slide image as a soft background if available
			if last_url:
				try:
					resp = requests.get(last_url, timeout=15)
					resp.raise_for_status()
					tmp_path = Path("_preview_thankyou_bg.png")
					tmp_path.write_bytes(resp.content)
					try:
						slide.shapes.add_picture(str(tmp_path), 0, 0, width=prs.slide_width, height=prs.slide_height)
						try:
							with Image.open(str(tmp_path)) as im:
								im_small = im.resize((64, 64)).convert("L")
								avg = sum(im_small.getdata()) / float(len(im_small.getdata()))
								if avg > 150:
									text_rgb = RGBColor(31, 41, 55)
								else:
									text_rgb = RGBColor(248, 250, 252)
						except Exception as sample_err:
							logger.debug(f"Failed to sample thank-you background brightness: {sample_err}")
					finally:
						try:
							tmp_path.unlink()
						except OSError:
							pass
				except Exception as e:
					logger.debug(f"Failed to add thank-you background image: {e}")

			from pptx.enum.text import PP_ALIGN

			# Centered "Thank you" text
			ty_box = slide.shapes.add_textbox(
				Inches(1.0),
				Inches(2.5),
				prs.slide_width - Inches(2.0),
				Inches(2.5),
			)
			ty_tf = ty_box.text_frame
			ty_tf.clear()
			ty_tf.word_wrap = True

			p_ty = ty_tf.paragraphs[0]
			p_ty.text = "Thank you"
			p_ty.alignment = PP_ALIGN.CENTER
			ty_font = p_ty.font
			ty_font.name = "Arial"
			ty_font.size = Pt(44)
			ty_font.bold = True
			ty_font.color.rgb = text_rgb

			# Optional subtitle
			p_sub = ty_tf.add_paragraph()
			p_sub.text = "Questions or discussion"
			p_sub.alignment = PP_ALIGN.CENTER
			sub_font = p_sub.font
			sub_font.name = "Arial"
			sub_font.size = Pt(20)
			sub_font.color.rgb = text_rgb

		# Filename similar to template-based export
		import re
		safe_title = re.sub(r"[^\w\s-]", "", title)[:50].strip().replace(" ", "_")
		if not safe_title:
			safe_title = "Presentation"

		safe_user = re.sub(r"[^\w\s-]", "", user_name)[:30].strip().replace(" ", "_")
		if not safe_user:
			safe_user = "user"

		date_str = datetime.utcnow().strftime("%Y-%m-%d")
		filename = f"{safe_title}_{safe_user}_{date_str}.pptx"

		return prs, filename

	def _build_presentation(self, deck_id: str, user_name: str = "user", export_style: str = "template") -> Tuple[Presentation, str]:
		# If user requests preview-style export, bypass templates entirely.
		if export_style == "preview":
			return self._build_preview_presentation(deck_id, user_name)
		try:
			object_id = ObjectId(deck_id)
		except Exception:
			raise ValueError("Invalid deck_id. Must be a Mongo ObjectId hex string.")

		deck = self.slides_collection.find_one({"_id": object_id})
		if not deck:
			raise FileNotFoundError("Slide deck not found")

		template_path = deck.get("template_path") or deck.get("metadata", {}).get("template_path")
		
		# Check if template is stored in MongoDB (prefixed with "db:")
		if template_path and template_path.startswith("db:"):
			# Load template from MongoDB
			template_id = template_path.replace("db:", "")
			templates_collection = self.db["templates"]
			template_doc = templates_collection.find_one({"templateId": template_id})
			
			if template_doc:
				try:
					storage_type = template_doc.get("storage_type", "base64")
					
					if storage_type == "gridfs":
						# Load from GridFS
						fs = GridFS(self.db, collection="templates_fs")
						gridfs_file_id = template_doc.get("template_file_id")
						
						if gridfs_file_id:
							gridfs_file = fs.get(ObjectId(gridfs_file_id))
							template_bytes = gridfs_file.read()
							prs = Presentation(io.BytesIO(template_bytes))
							logger.info(f"Loaded template '{template_id}' from GridFS ({len(template_bytes)} bytes)")
						else:
							raise ValueError("GridFS file ID not found")
					elif template_doc.get("template_file"):
						# Load from base64
						template_base64 = template_doc["template_file"]
						template_bytes = base64.b64decode(template_base64)
						prs = Presentation(io.BytesIO(template_bytes))
						logger.info(f"Loaded template '{template_id}' from MongoDB base64 ({len(template_bytes)} bytes)")
					else:
						raise ValueError("Template file data not found")
						
				except Exception as e:
					logger.warning(f"Failed to load template from DB: {e}, using default template")
					prs = Presentation()
			else:
				logger.warning(f"Template '{template_id}' not found in MongoDB, using default template")
				prs = Presentation()
		elif template_path and Path(template_path).exists():
			# Fallback to file system
			prs = Presentation(template_path)
		else:
			# No template specified, use default
			prs = Presentation()

		# Title slide – overwrite the first slide in the template if present,
		# otherwise add a new one using the title layout.
		if len(prs.slides):
			slide = prs.slides[0]
			title_shape = getattr(slide.shapes, "title", None)
			if title_shape is not None:
				title_shape.text = deck.get("title", "AI Presentation")
			else:
				# Fallback textbox if template has no explicit title shape
				box = slide.shapes.add_textbox(Inches(1.0), Inches(0.5), Inches(8.0), Inches(1.5))
				box.text = deck.get("title", "AI Presentation")
			try:
				subtitle = slide.placeholders[1]
				subtitle.text = f"Generated on {datetime.utcnow().strftime('%Y-%m-%d %H:%M UTC')}"
			except Exception:
				# Not all templates have a subtitle placeholder; ignore.
				pass
		else:
			title_layout = prs.slide_layouts[0]
			slide = prs.slides.add_slide(title_layout)
			if slide.shapes.title:
				slide.shapes.title.text = deck.get("title", "AI Presentation")
			try:
				subtitle = slide.placeholders[1]
				subtitle.text = f"Generated on {datetime.utcnow().strftime('%Y-%m-%d %H:%M UTC')}"
			except Exception:
				pass

		sections: List[str] = deck.get("sections", [])
		bullets: List[List[str]] = deck.get("bullets", [])
		expanded_content: List[str] = deck.get("expanded_content", [])  # Content after 3 passes
		generated_notes = deck.get("generated_notes", [])  # This contains expanded content from 3 passes
		speaker_notes = deck.get("speaker_notes", [])
		image_placeholders = deck.get("image_placeholders", [])
		media_refs = deck.get("media_refs", [])

		# Content slides
		for idx, section in enumerate(sections):
			# Try to reuse an existing template slide if available.
			target_index = idx + 1  # 0 is title slide
			if target_index < len(prs.slides):
				slide = prs.slides[target_index]
				title_shape = getattr(slide.shapes, "title", None)
				if title_shape is not None:
					title_shape.text = section
			else:
				content_layout = prs.slide_layouts[1]  # Title and Content
				slide = prs.slides.add_slide(content_layout)
				if slide.shapes.title:
					slide.shapes.title.text = section

			# Body text frame – fall back to a textbox if placeholder[1] is missing.
			used_placeholder = True
			try:
				text_frame = slide.placeholders[1].text_frame
			except Exception:
				used_placeholder = False
				box = slide.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(8.0), Inches(4.5))
				text_frame = box.text_frame

			# Capture the template's paragraph style before clearing so we can
			# re-use its font, size, spacing, etc. for all generated content.
			template_paragraph = text_frame.paragraphs[0] if text_frame.paragraphs else None

			# If we had to create our own textbox (no placeholder), try to borrow
			# style from the slide title so we still respect the template's font.
			if (template_paragraph is None) and (not used_placeholder):
				try:
					title_shape = getattr(slide.shapes, "title", None)
					if title_shape is not None and title_shape.text_frame.paragraphs:
						template_paragraph = title_shape.text_frame.paragraphs[0]
				except Exception:
					template_paragraph = None

			template_style = self._capture_paragraph_style(template_paragraph)

			# Extra safety net: if we still don't have a style (e.g. a very
			# minimal/clean template), synthesize a reasonable default based on
			# the slide title font so content still looks like it belongs to
			# the theme instead of falling back to a random system font.
			if not template_style:
				try:
					title_shape = getattr(slide.shapes, "title", None)
					title_paragraph = None
					if title_shape is not None and title_shape.text_frame.paragraphs:
						title_paragraph = title_shape.text_frame.paragraphs[0]
					if title_paragraph is not None:
						title_font = title_paragraph.font
						from pptx.util import Pt as _Pt
						template_style = {
							"font_name": title_font.name,
							# Use a slightly smaller size than the title if available,
							# otherwise fall back to a reasonable body size.
							"font_size": (
								title_font.size * 0.7 if title_font.size is not None else _Pt(24)
							),
							"bold": False,
							"italic": False,
							"font_color": getattr(getattr(title_font, "color", None), "rgb", None),
							"level": 0,
							"alignment": None,
							"space_after": None,
							"space_before": None,
						}
				except Exception:
					# If anything goes wrong here, we just keep template_style as-is
					# and let PowerPoint/theme defaults take over.
					pass

			# Now clear out any existing template text so we don't stack on top.
			text_frame.clear()
			
			# CRITICAL: Use expanded content (after 3 passes) instead of bullets
			# This is the actual expanded content from Gemini -> Gemini -> Deepseek
			# It's stored in generated_notes (which contains the expanded text after 3 passes)
			expanded_text = generated_notes[idx] if idx < len(generated_notes) else ""
			if not expanded_text and idx < len(expanded_content):
				expanded_text = expanded_content[idx]
			
			if expanded_text:
				# Place expanded content as paragraphs of text (not bullets)
				# Split into sentences for better formatting
				import re
				sentences = re.split(r'(?<=[.!?])\s+', expanded_text.strip())
				for i, sentence in enumerate(sentences):
					if sentence.strip():
						p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
						# Apply the original template styling so the new text looks
						# like it belongs to the slide design.
						self._apply_paragraph_style(p, template_style)
						p.text = sentence.strip()
			else:
				# Fallback to bullets if expanded content not available
				slide_bullets = bullets[idx] if idx < len(bullets) else []
				for i, bullet in enumerate(slide_bullets):
					p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
					self._apply_paragraph_style(p, template_style)
					p.text = bullet

			# Image placeholder captions (fallback text when no media available)
			placeholders = image_placeholders[idx] if idx < len(image_placeholders) else []
			for placeholder in placeholders:
				caption = placeholder.get("caption") or placeholder.get("marker") or placeholder.get("id")
				if not caption:
					continue
				p = text_frame.add_paragraph()
				p.text = f"[Image placeholder] {caption}"
				p.font.size = Pt(12)
				p.font.italic = True
				p.level = 1
				p.font.color.rgb = RGBColor(120, 120, 120)

			# Render stock / generated images when available
			slide_media = media_refs[idx] if idx < len(media_refs) else []
			if slide_media:
				# Use the first media URL for now; can be extended later.
				url = slide_media[0] if isinstance(slide_media[0], str) else None
				if url:
					try:
						response = requests.get(url, timeout=15)
						response.raise_for_status()
						img_bytes = response.content
						tmp_path = Path("_ppt_tmp_image.png")
						tmp_path.write_bytes(img_bytes)
						try:
							# Calculate available space dynamically based on text content
							available_space = self._calculate_available_space_for_image(slide, text_frame)
							# Calculate optimal image size
							optimal_size = self._calculate_optimal_image_size(str(tmp_path), available_space)
							
							# Add image with dynamically calculated size
							slide.shapes.add_picture(
								str(tmp_path),
								Inches(optimal_size["left"]),
								Inches(optimal_size["top"]),
								width=Inches(optimal_size["width"]),
								height=Inches(optimal_size["height"])
							)
						finally:
							try:
								tmp_path.unlink()
							except OSError:
								# Non-fatal if temp cleanup fails.
								pass
					except Exception as e:
						# If image download or placement fails, continue without blocking export.
						logger.debug(f"Failed to add image to slide {idx}: {e}")
						pass

			# Speaker notes priority
			notes_text = ""
			if idx < len(speaker_notes):
				note_entry = speaker_notes[idx]
				main_points = note_entry.get("main_points") or []
				talking_points = note_entry.get("talking_points") or []
				if main_points or talking_points:
					chunks = []
					if main_points:
						chunks.append("Main points: " + "; ".join(main_points[:4]))
					if talking_points:
						chunks.append("Talking points: " + "; ".join(talking_points[:4]))
					notes_text = "\n".join(chunks)
			if not notes_text and idx < len(generated_notes):
				notes_text = generated_notes[idx] or ""

			if notes_text:
				notes_frame = slide.notes_slide.notes_text_frame
				notes_frame.clear()
				notes_frame.text = notes_text.strip()

		# Smart slide removal: Keep important conclusion slides (Thank You, Questions, etc.) and remove middle slides
		# Strategy: Keep title (0) + user content (1 to len(sections)) + important end slides (last 1-2)
		def _is_important_conclusion_slide(slide) -> bool:
			"""Detect if slide is an important conclusion slide (Thank You, Questions, etc.)"""
			conclusion_keywords = [
				"thank you", "thanks", "thank", "questions", "q&a", "qa", 
				"conclusion", "conclusions", "summary", "contact", "reach out",
				"any questions", "questions?", "thank", "appreciation"
			]
			try:
				# Check title
				if slide.shapes.title and slide.shapes.title.text:
					title_text = slide.shapes.title.text.lower().strip()
					if any(keyword in title_text for keyword in conclusion_keywords):
						return True
				# Check all text shapes
				for shape in slide.shapes:
					if hasattr(shape, "text") and shape.text:
						text = shape.text.lower().strip()
						if any(keyword in text for keyword in conclusion_keywords):
							return True
			except:
				pass
			return False
		
		# Detect important conclusion slides (usually last 1-2 slides)
		important_end_slides = []
		for i in range(max(0, len(prs.slides) - 2), len(prs.slides)):
			if i > 0 and _is_important_conclusion_slide(prs.slides[i]):
				important_end_slides.append(i)
		
		# Calculate what we need: title (1) + user content (len(sections)) + important end slides
		num_important_end = len(important_end_slides)
		total_needed = 1 + len(sections) + num_important_end
		
		if len(prs.slides) > total_needed:
			slides_to_remove = len(prs.slides) - total_needed
			logger.info(f"Template has {len(prs.slides)} slides, need {total_needed} (title + {len(sections)} content + {num_important_end} conclusion). Removing {slides_to_remove} middle slide(s).")
			
			# Remove middle slides, keeping: title (0), user content area, and important end slides
			# We'll remove slides from the middle range, not from the end
			try:
				sldIdLst = prs.slides._sldIdLst
				# Calculate which slides to remove (middle ones, not end ones)
				# Keep: slide 0 (title), slides 1 to len(sections) (user content), and important end slides
				slides_to_keep_indices = set([0])  # Title slide
				slides_to_keep_indices.update(range(1, 1 + len(sections)))  # User content slides
				slides_to_keep_indices.update(important_end_slides)  # Important conclusion slides
				
				# Remove slides that are NOT in the keep list (work backwards to avoid index issues)
				indices_to_remove = []
				for i in range(len(prs.slides) - 1, -1, -1):
					if i not in slides_to_keep_indices:
						indices_to_remove.append(i)
				
				# Remove slide references (work backwards)
				for idx in sorted(indices_to_remove, reverse=True):
					if idx < len(sldIdLst):
						try:
							sldIdLst.remove(sldIdLst[idx])
						except:
							pass
				
				logger.info(f"Successfully removed {len(indices_to_remove)} middle slide(s), kept {num_important_end} important conclusion slide(s).")
			except Exception as e:
				# If removal fails, at least clear the content of middle slides (keep end slides intact)
				logger.warning(f"Could not remove slide references: {e}. Clearing content of middle slides instead.")
				slides_to_keep_indices = set([0])
				slides_to_keep_indices.update(range(1, 1 + len(sections)))
				slides_to_keep_indices.update(important_end_slides)
				
				for i in range(len(prs.slides)):
					if i not in slides_to_keep_indices:
						try:
							extra_slide = prs.slides[i]
							# Clear all shapes on the extra slide
							for shape in list(extra_slide.shapes):
								try:
									sp = shape._element
									sp.getparent().remove(sp)
								except:
									pass
						except Exception as e2:
							logger.debug(f"Could not clear slide {i}: {e2}")

		# Generate filename: Topic_user_date.pptx
		title = deck.get("title", "Presentation")
		# Sanitize title for filename (remove special chars, limit length)
		import re
		safe_title = re.sub(r'[^\w\s-]', '', title)[:50].strip().replace(' ', '_')
		if not safe_title:
			safe_title = "Presentation"
		
		# Sanitize user name
		safe_user = re.sub(r'[^\w\s-]', '', user_name)[:30].strip().replace(' ', '_')
		if not safe_user:
			safe_user = "user"
		
		# Format date as YYYY-MM-DD
		date_str = datetime.utcnow().strftime('%Y-%m-%d')
		
		filename = f"{safe_title}_{safe_user}_{date_str}.pptx"
		
		return prs, filename
