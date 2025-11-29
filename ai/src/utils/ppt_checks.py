from __future__ import annotations

from typing import List, Tuple, Union
from io import BytesIO

from pptx import Presentation

FORBIDDEN_TOKENS = ["```", "{", "}", '"slides":', "Notes:"]

PptSource = Union[str, bytes, bytearray]


def _load_presentation(source: PptSource) -> Presentation:
	if isinstance(source, (bytes, bytearray)):
		return Presentation(BytesIO(source))
	return Presentation(source)


def check_no_json_tokens(ppt_source: PptSource) -> List[Tuple[int, str, str]]:
	"""Return a list of (slide_index, token, snippet) for offending text."""
	prs = _load_presentation(ppt_source)
	failures: List[Tuple[int, str, str]] = []
	for slide_idx, slide in enumerate(prs.slides, start=1):
		for shape in slide.shapes:
			if not getattr(shape, "has_text_frame", False):
				continue
			text = shape.text_frame.text
			if not text:
				continue
			for token in FORBIDDEN_TOKENS:
				if token in text:
					failures.append((slide_idx, token, text[:80]))
	return failures


def check_bullets_limit(ppt_source: PptSource, max_bullets: int = 6) -> List[Tuple[int, int]]:
	"""Return slides that exceed the bullet ceiling."""
	prs = _load_presentation(ppt_source)
	failures: List[Tuple[int, int]] = []
	for slide_idx, slide in enumerate(prs.slides, start=1):
		bullet_count = 0
		for shape in slide.shapes:
			if not getattr(shape, "has_text_frame", False):
				continue
			text_frame = shape.text_frame
			bullet_count += sum(1 for paragraph in text_frame.paragraphs if paragraph.text.strip())
		if bullet_count > max_bullets + 2:
			failures.append((slide_idx, bullet_count))
	return failures

