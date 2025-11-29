import os
import tempfile
import unittest

from pptx import Presentation

from utils.slide_cleaner import (
	clean_text,
	sanitize_slide_payload,
	SlideValidationError,
)
from utils.ppt_checks import check_bullets_limit, check_no_json_tokens
from utils.fix_ppt_pipeline import prepare_slides_from_raw


class SlideCleanerTests(unittest.TestCase):
	def test_clean_text_strips_tokens(self):
		text = "Notes: ```json\n{\"key\": \"value\"}\n```"
		self.assertEqual(clean_text(text), '"key": "value"')

	def test_sanitize_limits_bullets(self):
		payload = {
			"slides": [{
				"title": "Test",
				"bullets": [
					"First bullet that is intentionally long to be truncated at some point",
					"",
					"Second bullet"
				],
				"notes": ""
			}]
		}
		sanitized = sanitize_slide_payload(payload)
		self.assertEqual(len(sanitized["slides"][0]["bullets"]), 2)
		self.assertTrue(all(len(b.split()) <= 12 for b in sanitized["slides"][0]["bullets"]))

	def test_missing_slides_raises(self):
		with self.assertRaises(SlideValidationError):
			sanitize_slide_payload({})


class PptChecksTests(unittest.TestCase):
	def _make_ppt(self, body_text: str) -> bytes:
		prs = Presentation()
		slide = prs.slides.add_slide(prs.slide_layouts[1])
		slide.shapes.title.text = "Test"
		text_frame = slide.shapes.placeholders[1].text_frame
		text_frame.text = body_text
		buffer = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
		prs.save(buffer.name)
		with open(buffer.name, "rb") as fh:
			data = fh.read()
		self.addCleanup(lambda: os.remove(buffer.name) if os.path.exists(buffer.name) else None)
		return data

	def test_check_no_json_tokens_detects(self):
		path = self._make_ppt("Contains ``` token")
		failures = check_no_json_tokens(path)
		self.assertTrue(failures)

	def test_check_bullets_limit_detects(self):
		path = self._make_ppt("Line1\nLine2\nLine3\nLine4\nLine5\nLine6\nLine7")
		failures = check_bullets_limit(path, max_bullets=3)
		self.assertTrue(failures)


class FixPipelineTests(unittest.TestCase):
	def test_prepare_slides_from_raw_fills_missing(self):
		raw = """
Noise...
{
  "meta": {"presentation_title": "Photosynthesis deck"},
  "slides": [
    {
      "title": "Intro",
      "bullets": ["Explain process", "Mention chlorophyll"],
      "notes": "Notes: keep it simple"
    }
  ]
}
"""
		prepared = prepare_slides_from_raw(raw, desired_slide_count=3)
		self.assertEqual(prepared["meta"]["presentation_title"], "Photosynthesis deck")
		self.assertEqual(len(prepared["slides"]), 3)
		self.assertFalse(prepared["slides"][0]["notes"].lower().startswith("notes"))


if __name__ == "__main__":
	unittest.main()

