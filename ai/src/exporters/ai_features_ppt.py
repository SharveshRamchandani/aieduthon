from datetime import datetime
from pathlib import Path
from typing import List, Dict

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

FEATURES: List[Dict[str, List[str]]] = [
	{
		"title": "1. Prompt-to-Slide Generation",
		"bullets": [
			"Converts text prompts into full, sequenced slides",
			"Structures content and examples per topic",
			"Novelty: Fully automates deck building; saves teacher hours"
		]
	},
	{
		"title": "2. Smart Quiz Injection",
		"bullets": [
			"Extracts key facts using NLP from slides",
			"Creates MCQs and interactive quizzes per topic",
			"Novelty: Context-embedded quizzes; exportable for live use"
		]
	},
	{
		"title": "3. Speaker Notes Automation",
		"bullets": [
			"Generates detailed speaker notes aligned to slides",
			"Adapts to audience and tone",
			"Novelty: In-slide cues boost presentation fluency"
		]
	},
	{
		"title": "4. Concept-to-Diagram Generator",
		"bullets": [
			"Builds educational diagrams via NLP + Python viz",
			"Flowcharts, processes, systems",
			"Novelty: Fills visual learning gap in slide generators"
		]
	},
	{
		"title": "5. Cultural & Linguistic Localization",
		"bullets": [
			"Adapts language, examples, and images to locale",
			"Uses profiles/context for hyper-localization",
			"Novelty: Beyond translation; India-focused relevance"
		]
	},
	{
		"title": "6. Media Integration",
		"bullets": [
			"Fetches/generates images, diagrams, infographics",
			"NLP-based tagging, open-source APIs",
			"Novelty: Optimized by subject, region, audience"
		]
	},
	{
		"title": "7. Template/Style Auto-Selection",
		"bullets": [
			"Chooses best-fit templates/styles",
			"Considers topic, teacher preference, audience",
			"Novelty: Customizes look and accessibility"
		]
	},
	{
		"title": "8. Multi-language Output",
		"bullets": [
			"Generates slides/quizzes/notes in Indian languages",
			"No manual formatting needed",
			"Novelty: Empowers multilingual teaching"
		]
	},
	{
		"title": "9. Analytics & Feedback Agent",
		"bullets": [
			"Tracks time saved, popularity, quiz accuracy, styles",
			"Dashboards for admins/teachers",
			"Novelty: Measures real educational impact"
		]
	},
	{
		"title": "10. Offline-Friendly Generation",
		"bullets": [
			"Supports low/zero-bandwidth slide/quiz/diagram creation",
			"Caching and deferred sync",
			"Novelty: Access in rural and constrained classrooms"
		]
	},
	{
		"title": "11. Agentic Orchestration",
		"bullets": [
			"End-to-end pipeline with context passing",
			"Improves via analytics and feedback",
			"Novelty: Self-coordinating intelligent workflow"
		]
	}
]


def build_ai_features_ppt(output_dir: str = "..\\..\\out", filename: str = "ai_features.pptx") -> str:
	prs = Presentation()

	# Title slide
	title_slide = prs.slides.add_slide(prs.slide_layouts[0])
	title_slide.shapes.title.text = "Complete AI Features"
	sub = title_slide.placeholders[1]
	sub.text = f"Auto-generated on {datetime.utcnow().strftime('%Y-%m-%d %H:%M UTC')}"

	# Slides for each feature
	for feature in FEATURES:
		slide = prs.slides.add_slide(prs.slide_layouts[1])
		slide.shapes.title.text = feature["title"]
		frame = slide.shapes.placeholders[1].text_frame
		frame.clear()
		for b in feature["bullets"]:
			p = frame.add_paragraph()
			p.text = b
			p.font.size = Pt(20)

	# Save
	out = Path(output_dir).resolve()
	out.mkdir(parents=True, exist_ok=True)
	path = out / filename
	prs.save(str(path))
	return str(path)
