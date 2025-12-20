import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# Words that strongly indicate lorem ipsum / placeholder text
LOREM_MARKERS = {
    "lorem",
    "ipsum",
    "dolor",
    "sit",
    "amet",
    "consectetur",
    "adipiscing",
    "elit",
    "dummy",
    "placeholder",
    "sample text",
    "your title here",
    "your subtitle",
    "click to add title",
    "click to add text",
}

# If True, remove ALL images from templates.
# If False, only remove images that look like placeholders by name/alt text.
REMOVE_ALL_IMAGES = False


def is_lorem_text(text: str) -> bool:
    """Heuristic: decide if a text string is mostly lorem/placeholder."""
    if not text:
        return False

    lower = text.lower()
    # Quick substring check
    if any(marker in lower for marker in LOREM_MARKERS):
        return True

    # If it's very short generic text, also consider it placeholder-ish
    stripped = lower.strip()
    if stripped in {"title", "subtitle", "heading", "subheading"}:
        return True

    return False


def clean_text_shapes(prs: Presentation) -> int:
    """Remove ALL text from all shapes in the presentation, keep shapes/styles."""
    removed_count = 0

    for slide in prs.slides:
        for shape in list(slide.shapes):
            # Text in standard shapes / placeholders
            if hasattr(shape, "text_frame") and shape.text_frame is not None:
                if shape.text_frame.text:
                    shape.text_frame.text = ""
                    removed_count += 1
            elif hasattr(shape, "text") and isinstance(shape.text, str):
                if shape.text:
                    try:
                        shape.text = ""
                        removed_count += 1
                    except Exception:
                        pass

    return removed_count


def looks_like_placeholder_image(shape) -> bool:
    """Heuristic for detecting placeholder/sample images."""
    name = getattr(shape, "name", "") or ""
    alt_text = getattr(shape, "alternative_text", "") or ""

    combined = f"{name} {alt_text}".lower()
    placeholder_markers = [
        "placeholder",
        "sample",
        "lorem",
        "dummy",
        "example",
        "your image here",
        "drag image here",
        "stock photo",
    ]
    return any(m in combined for m in placeholder_markers)


def clean_images(prs: Presentation) -> int:
    """Remove ALL images from the presentation, keep layout shapes."""
    removed = 0

    for slide in prs.slides:
        # work on a copy to safely remove shapes
        for shape in list(slide.shapes):
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    sp = shape._element
                    sp.getparent().remove(sp)
                    removed += 1
                except Exception:
                    # Non-fatal if we can't remove a specific image
                    pass

    return removed


def clean_pptx_file(path: Path, backup: bool = False) -> None:
    print(f"Cleaning: {path}")
    prs = Presentation(str(path))

    removed_text = clean_text_shapes(prs)
    removed_images = clean_images(prs)

    # By default we now overwrite in place to avoid Windows file-lock issues
    # when attempting to rename/open files. If you want backups, run this
    # script on a copy of the templates folder.
    prs.save(str(path))
    print(f"  Removed lorem/placeholder text from {removed_text} shape(s)")
    print(f"  Removed {removed_images} image(s)")
    print("  Saved cleaned file.\n")


def collect_pptx_files(target: Path):
    if target.is_file() and target.suffix.lower() == ".pptx":
        return [target]
    elif target.is_dir():
        return sorted(target.rglob("*.pptx"))
    else:
        return []


def main():
    if len(sys.argv) < 2:
        print("Usage: python clean_ppt_templates.py <pptx_file_or_folder>")
        sys.exit(1)

    target = Path(sys.argv[1]).resolve()
    files = collect_pptx_files(target)

    if not files:
        print(f"No .pptx files found at {target}")
        sys.exit(1)

    print(f"Found {len(files)} PPTX file(s) to clean.\n")
    for f in files:
        clean_pptx_file(f, backup=False)


if __name__ == "__main__":
    main()


