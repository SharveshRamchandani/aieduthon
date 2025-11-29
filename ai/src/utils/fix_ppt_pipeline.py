"""
Utility helpers to recover noisy LLM slide output, validate it, generate images,
and build a clean PPTX deck fully populated with notes and visuals.
"""

from __future__ import annotations

import io
import json
import logging
import os
import re
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

from jsonschema import validate, ValidationError
from pptx import Presentation
from pptx.util import Inches
from PIL import Image

from agents.image_generation_agent import ImageGenerationAgent
from utils.slide_cleaner import clean_text, sanitize_slide_payload, SlideValidationError

logger = logging.getLogger(__name__)

IMAGES_DIR = Path("ppt_images")
IMAGES_DIR.mkdir(parents=True, exist_ok=True)

SCHEMA: Dict[str, Any] = {
    "type": "object",
    "properties": {
        "meta": {"type": "object"},
        "slides": {"type": "array"},
    },
    "required": ["slides"],
}

TEMPLATE_MAP = {
    "title": 0,
    "title_content": 1,
    "bullet_list": 1,
    "two_column": 3,
    "image_fullbleed": 6,
    "qa": 1,
}

DEFAULT_IMAGE_BOX = {"left": 6.2, "top": 1.0, "w": 3.5, "h": 4.5}
DEFAULT_DPI = 150
_IMAGE_AGENT: Optional[ImageGenerationAgent] = None


def _get_image_agent() -> ImageGenerationAgent:
    global _IMAGE_AGENT
    if _IMAGE_AGENT is None:
        _IMAGE_AGENT = ImageGenerationAgent()
    return _IMAGE_AGENT


def extract_json_from_text(text: str) -> Dict[str, Any]:
    text = (text or "").strip()
    if not text:
        raise ValueError("LLM output is empty")
    if text.startswith("{"):
        try:
            return json.loads(text)
        except json.JSONDecodeError:
            pass
    stack: List[int] = []
    start: Optional[int] = None
    for idx, ch in enumerate(text):
        if ch == "{":
            if start is None:
                start = idx
            stack.append(idx)
        elif ch == "}":
            if stack:
                stack.pop()
                if not stack and start is not None:
                    candidate = text[start : idx + 1]
                    try:
                        return json.loads(candidate)
                    except json.JSONDecodeError:
                        start = None
    raise ValueError("Could not recover JSON payload from LLM text")


def call_image_api(prompt: str, width: int, height: int, seed: int) -> Tuple[bytes, str]:
    agent = _get_image_agent()
    metadata = {"seed": seed, "description": prompt}
    result = agent.generate(
        prompt=prompt,
        width=width,
        height=height,
        metadata=metadata,
        use_cache=True,
    )
    if not result.get("success"):
        raise RuntimeError(result.get("error", "Image generation failed"))
    url = (result.get("urls") or [None])[0]
    if not url:
        raise RuntimeError("Image generation returned no URL")
    filename = url.split("/")[-1]
    local_path = Path("out/generated_images") / filename
    if not local_path.exists():
        raise FileNotFoundError(f"Generated image file missing: {local_path}")
    with open(local_path, "rb") as handle:
        return handle.read(), str(local_path)


def _crop_center_to_aspect(img: Image.Image, target_w: int, target_h: int) -> Image.Image:
    iw, ih = img.size
    target_ratio = target_w / target_h
    image_ratio = iw / ih
    if image_ratio > target_ratio:
        new_w = int(ih * target_ratio)
        left = (iw - new_w) // 2
        box = (left, 0, left + new_w, ih)
    else:
        new_h = int(iw / target_ratio)
        top = (ih - new_h) // 2
        box = (0, top, iw, top + new_h)
    return img.crop(box)


def _add_image(slide, img_path: str, box: Dict[str, float], mode: str = "fill"):
    image = Image.open(img_path)
    box_left, box_top, box_w, box_h = (
        box.get("left", DEFAULT_IMAGE_BOX["left"]),
        box.get("top", DEFAULT_IMAGE_BOX["top"]),
        box.get("w", DEFAULT_IMAGE_BOX["w"]),
        box.get("h", DEFAULT_IMAGE_BOX["h"]),
    )
    if mode == "fill":
        cropped = _crop_center_to_aspect(
            image,
            int(box_w * DEFAULT_DPI),
            int(box_h * DEFAULT_DPI),
        )
        tmp_path = Path(img_path).with_suffix(".crop.png")
        cropped.save(tmp_path)
        slide.shapes.add_picture(
            str(tmp_path),
            Inches(box_left),
            Inches(box_top),
            width=Inches(box_w),
            height=Inches(box_h),
        )
        tmp_path.unlink(missing_ok=True)
    else:
        iw, ih = image.size
        img_ratio = iw / ih
        box_ratio = box_w / box_h
        if img_ratio > box_ratio:
            used_w = box_w
            used_h = used_w / img_ratio
        else:
            used_h = box_h
            used_w = used_h * img_ratio
        left = box_left + (box_w - used_w) / 2
        top = box_top + (box_h - used_h) / 2
        slide.shapes.add_picture(
            str(Path(img_path)),
            Inches(left),
            Inches(top),
            width=Inches(used_w),
            height=Inches(used_h),
        )


def _build_slide_prompt(slide: Dict[str, Any]) -> str:
    caption = None
    images = slide.get("images") or []
    if images:
        entry = images[0]
        caption = entry.get("caption") or entry.get("marker")
    if not caption:
        caption = slide.get("title")
    bullets = slide.get("bullets") or []
    summary = " ".join(bullets[:3])
    return clean_text(f"{caption}. {summary}")


def _ensure_image(prompt: str, idx: int, box: Optional[Dict[str, float]], seed_base: int) -> Optional[str]:
    try:
        width = int(round((box or DEFAULT_IMAGE_BOX)["w"] * DEFAULT_DPI))
        height = int(round((box or DEFAULT_IMAGE_BOX)["h"] * DEFAULT_DPI))
        bytes_data, _ = call_image_api(prompt, width, height, seed_base + idx)
        img_path = IMAGES_DIR / f"slide_{idx}.png"
        with open(img_path, "wb") as handle:
            handle.write(bytes_data)
        return str(img_path)
    except Exception as exc:
        logger.warning("Image generation failed for slide %s: %s", idx, exc)
        return None


def _fallback_slide(index: int) -> Dict[str, Any]:
    title = f"Topic {index}"
    bullets = [
        "Key concept overview",
        "Supporting detail or fact",
        "Classroom example"
    ]
    return {
        "title": title,
        "bullets": bullets,
        "notes": "",
        "template": "title_content",
        "image_prompt": title,
        "image_box_inches": DEFAULT_IMAGE_BOX,
        "image_mode": "fill",
        "images": [{"caption": title}],
    }


def prepare_slides_from_raw(
    raw_text: str,
    desired_slide_count: Optional[int] = None,
) -> Dict[str, Any]:
    payload = extract_json_from_text(raw_text)
    try:
        validate(payload, SCHEMA)
    except ValidationError as exc:
        logger.warning("Slide JSON failed validation: %s", exc)
    slides_raw = payload.get("slides", []) or []
    meta = payload.get("meta") or {}
    cleaned_slides: List[Dict[str, Any]] = []
    try:
        sanitized = sanitize_slide_payload({"slides": slides_raw})
        cleaned_slides = sanitized.get("slides", [])
    except SlideValidationError as exc:
        logger.warning("Sanitizer rejected slide payload: %s", exc)
        cleaned_slides = []
    normalized_slides: List[Dict[str, Any]] = []
    for idx, raw_slide in enumerate(slides_raw):
        clean_slide = cleaned_slides[idx] if idx < len(cleaned_slides) else {}
        title = clean_slide.get("title") or clean_text(raw_slide.get("title")) or f"Slide {idx + 1}"
        bullets = clean_slide.get("bullets") or [clean_text(b) for b in (raw_slide.get("bullets") or []) if clean_text(b)]
        if not bullets:
            bullets = _fallback_slide(idx + 1)["bullets"]
        notes = clean_slide.get("notes") or clean_text(raw_slide.get("notes") or "")
        images = clean_slide.get("images") or raw_slide.get("images") or [{"caption": title}]
        normalized_slides.append({
            "title": title,
            "bullets": bullets,
            "notes": notes,
            "template": raw_slide.get("template", "title_content"),
            "images": images,
            "image_prompt": clean_text(raw_slide.get("image_prompt") or images[0].get("caption") or title),
            "image_box_inches": raw_slide.get("image_box_inches") or DEFAULT_IMAGE_BOX,
            "image_mode": raw_slide.get("image_mode", "fill"),
        })
    target_count = desired_slide_count or len(normalized_slides)
    while len(normalized_slides) < target_count:
        normalized_slides.append(_fallback_slide(len(normalized_slides) + 1))
    if not normalized_slides:
        normalized_slides.append(_fallback_slide(1))
    return {"meta": meta, "slides": normalized_slides}


def build_clean_ppt_from_raw(
    raw_text: str,
    deck_title: Optional[str] = None,
    seed_base: int = 1000,
    template_path: Optional[str] = None,
) -> Tuple[bytes, str]:
    prepared = prepare_slides_from_raw(raw_text)
    slides = prepared["slides"]
    meta = prepared["meta"]
    if template_path and Path(template_path).exists():
        prs = Presentation(template_path)
    else:
        prs = Presentation()
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = meta.get("presentation_title") or deck_title or "AI Presentation"
    subtitle = title_slide.placeholders[1]
    subtitle.text = meta.get("subtitle") or "Generated automatically"
    for idx, slide_data in enumerate(slides, start=1):
        template = slide_data.get("template", "title_content")
        layout_index = TEMPLATE_MAP.get(template, 1)
        try:
            layout = prs.slide_layouts[layout_index]
        except IndexError:
            layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)
        title = slide_data.get("title")
        if title:
            try:
                slide.shapes.title.text = clean_text(title)
            except Exception:
                box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9.0), Inches(1.0))
                box.text = clean_text(title)
        bullets = slide_data.get("bullets") or []
        if bullets:
            try:
                text_frame = slide.placeholders[1].text_frame
                text_frame.clear()
            except Exception:
                box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(8.5), Inches(4.5))
                text_frame = box.text_frame
                text_frame.clear()
            for i, bullet in enumerate(bullets):
                paragraph = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
                paragraph.text = clean_text(bullet)
        notes = slide_data.get("notes")
        if notes:
            slide.notes_slide.notes_text_frame.text = clean_text(notes)
        image_prompt = slide_data.get("image_prompt") or _build_slide_prompt(slide_data)
        box = slide_data.get("image_box_inches") or DEFAULT_IMAGE_BOX
        mode = slide_data.get("image_mode", meta.get("default_image_mode", "fill"))
        img_path = _ensure_image(image_prompt, idx, box, seed_base)
        if img_path:
            _add_image(slide, img_path, box, mode=mode)
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    filename = f"{(deck_title or meta.get('presentation_title') or 'deck').replace(' ', '_')}.pptx"
    return buffer.getvalue(), filename


def main():
    raw_path = Path("llm_raw.txt")
    if not raw_path.exists():
        raise SystemExit("llm_raw.txt not found")
    with raw_path.open("r", encoding="utf-8") as file:
        raw_text = file.read()
    ppt_bytes, fname = build_clean_ppt_from_raw(raw_text)
    output = Path("fixed_presentation_final.pptx")
    output.write_bytes(ppt_bytes)
    logger.info("Saved cleaned PPT to %s (suggested filename: %s)", output, fname)


if __name__ == "__main__":
    main()

